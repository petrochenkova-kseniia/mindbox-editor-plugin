#!/usr/bin/env python3
"""
Mindbox Presentation Builder

Клонирует слайды из шаблона Mindbox и заменяет текст.

Использование:
    python3 pptx_builder.py build <config.json> <output.pptx>
    python3 pptx_builder.py list-slides
    python3 pptx_builder.py inspect <slide_num> [slide_num2 ...]

Формат config.json:
{
  "slides": [
    {
      "template_slide": 4,
      "replacements": {
        "Заголовок слайда": "Мой заголовок"
      }
    }
  ]
}
"""

import sys
import json
import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATE_PATH = Path(__file__).parent / "template.pptx"


def _normalize_text(text):
    """Нормализовать спецсимволы для надёжного матчинга."""
    text = text.replace('\xa0', ' ')   # неразрывный пробел → обычный
    text = text.replace('\x0b', '')    # вертикальный таб → убрать
    text = text.replace('\r', '')      # carriage return → убрать
    return text.strip()


NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_CT = 'http://schemas.openxmlformats.org/package/2006/content-types'
NS_REL = 'http://schemas.openxmlformats.org/package/2006/relationships'


def build_presentation(config, output_path, template_path=None):
    """
    Собрать презентацию из конфига через ZIP-манипуляцию.

    Подход:
    1. Распаковать шаблон
    2. Скопировать нужные слайды (с дубликатами)
    3. Обновить presentation.xml и [Content_Types].xml
    4. Запаковать обратно
    5. Открыть через python-pptx для замены текста
    """
    if template_path is None:
        template_path = TEMPLATE_PATH
    template_path = Path(template_path)

    slides_config = config.get("slides", [])
    if not slides_config:
        raise ValueError("Конфиг не содержит слайдов")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)

        # Распаковать шаблон
        src_dir = tmpdir / "src"
        with zipfile.ZipFile(str(template_path), 'r') as z:
            z.extractall(str(src_dir))

        # Определяем маппинг template_slide (1-based) -> файлы
        # Слайды в PPTX: ppt/slides/slide1.xml, slide2.xml, ...
        # Но порядок определяется в ppt/presentation.xml через sldIdLst

        # Парсим presentation.xml чтобы узнать порядок слайдов
        pres_xml_path = src_dir / "ppt" / "presentation.xml"
        pres_tree = etree.parse(str(pres_xml_path))
        pres_root = pres_tree.getroot()

        sldIdLst = pres_root.find(f'{{{NS_P}}}sldIdLst')
        sld_ids = list(sldIdLst)

        # Парсим rels чтобы rId -> slide file
        pres_rels_path = src_dir / "ppt" / "_rels" / "presentation.xml.rels"
        rels_tree = etree.parse(str(pres_rels_path))
        rels_root = rels_tree.getroot()

        rid_to_target = {}
        for rel in rels_root:
            rid_to_target[rel.get('Id')] = rel.get('Target')

        # Строим маппинг: slide_index (1-based) -> slide filename
        slide_order = []  # list of (rId, target_file)
        for sld_id in sld_ids:
            rId = sld_id.get(f'{{{NS_R}}}id')
            target = rid_to_target.get(rId, '')
            slide_order.append((rId, target))

        # Теперь собираем output
        out_dir = tmpdir / "out"
        shutil.copytree(str(src_dir), str(out_dir))

        # Удаляем все слайды из out
        out_slides_dir = out_dir / "ppt" / "slides"
        out_slides_rels_dir = out_dir / "ppt" / "slides" / "_rels"

        for f in out_slides_dir.glob("slide*.xml"):
            f.unlink()
        if out_slides_rels_dir.exists():
            for f in out_slides_rels_dir.glob("slide*.xml.rels"):
                f.unlink()

        # Копируем нужные слайды с переименованием
        new_rids = []
        for i, sc in enumerate(slides_config):
            src_slide_idx = sc["template_slide"] - 1  # 0-based
            if src_slide_idx >= len(slide_order):
                raise ValueError(f"Слайд {sc['template_slide']} не существует в шаблоне (макс: {len(slide_order)})")

            _, src_target = slide_order[src_slide_idx]
            src_filename = os.path.basename(src_target)  # e.g. "slide4.xml"

            new_filename = f"slide{i + 1}.xml"
            new_rid = f"rIdSlide{i + 1}"

            # Копируем XML слайда
            src_slide_path = src_dir / "ppt" / "slides" / src_filename
            dst_slide_path = out_dir / "ppt" / "slides" / new_filename
            shutil.copy2(str(src_slide_path), str(dst_slide_path))

            # Копируем rels слайда
            src_rels_path = src_dir / "ppt" / "slides" / "_rels" / f"{src_filename}.rels"
            if src_rels_path.exists():
                dst_rels_path = out_dir / "ppt" / "slides" / "_rels" / f"{new_filename}.rels"
                shutil.copy2(str(src_rels_path), str(dst_rels_path))

            new_rids.append((new_rid, f"slides/{new_filename}"))

        # Обновляем presentation.xml
        # Удаляем все sldId
        for child in list(sldIdLst):
            sldIdLst.remove(child)

        # Добавляем новые
        for i, (new_rid, _) in enumerate(new_rids):
            sld_id_elem = etree.SubElement(sldIdLst, f'{{{NS_P}}}sldId')
            sld_id_elem.set('id', str(256 + i))
            sld_id_elem.set(f'{{{NS_R}}}id', new_rid)

        pres_tree.write(str(out_dir / "ppt" / "presentation.xml"),
                       xml_declaration=True, encoding='UTF-8', standalone=True)

        # Обновляем presentation.xml.rels
        # Удаляем старые slide relationships, добавляем новые
        slide_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"

        to_remove = []
        for rel in rels_root:
            if rel.get('Type') == slide_rel_type:
                to_remove.append(rel)
        for rel in to_remove:
            rels_root.remove(rel)

        for new_rid, target in new_rids:
            rel_elem = etree.SubElement(rels_root, 'Relationship')
            rel_elem.set('Id', new_rid)
            rel_elem.set('Type', slide_rel_type)
            rel_elem.set('Target', target)

        rels_tree.write(str(out_dir / "ppt" / "_rels" / "presentation.xml.rels"),
                       xml_declaration=True, encoding='UTF-8', standalone=True)

        # Обновляем [Content_Types].xml
        ct_path = out_dir / "[Content_Types].xml"
        ct_tree = etree.parse(str(ct_path))
        ct_root = ct_tree.getroot()

        # Удаляем старые Override для слайдов
        to_remove = []
        for override in ct_root:
            pn = override.get('PartName', '')
            if re.match(r'/ppt/slides/slide\d+\.xml$', pn):
                to_remove.append(override)
        for elem in to_remove:
            ct_root.remove(elem)

        # Добавляем новые
        slide_ct = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
        for i in range(len(slides_config)):
            override = etree.SubElement(ct_root, 'Override')
            override.set('PartName', f'/ppt/slides/slide{i + 1}.xml')
            override.set('ContentType', slide_ct)

        ct_tree.write(str(ct_path), xml_declaration=True, encoding='UTF-8', standalone=True)

        # Запаковываем обратно в PPTX
        raw_pptx = tmpdir / "raw_output.pptx"
        _zip_directory(out_dir, raw_pptx)

        # Открываем через python-pptx для замены текста
        prs = Presentation(str(raw_pptx))

        replacements_count = 0
        cleared_count = 0
        deleted_images = 0

        for i, sc in enumerate(slides_config):
            if i >= len(prs.slides):
                print(f"⚠️ Слайд {i+1} (template:{sc.get('template_slide')}) — ПРОПУЩЕН, нет в презентации!")
                continue

            slide = prs.slides[i]
            replacements = sc.get("replacements", {})
            image_placeholder = sc.get("image_placeholder")

            # 1. Заменяем текст / очищаем шейпы без замен
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                shape_text = _normalize_text(shape.text)
                if not shape_text:
                    continue

                matched = False
                for old_text, new_text in replacements.items():
                    if new_text == "":
                        new_text = " "
                    old_stripped = _normalize_text(old_text)
                    if shape_text == old_stripped:
                        _replace_all_text(shape, new_text)
                        matched = True
                        replacements_count += 1
                        break
                    if old_stripped in shape_text:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if old_stripped in run.text:
                                    run.text = run.text.replace(old_stripped, new_text)
                        matched = True
                        replacements_count += 1
                        break

                if not matched:
                    # Очищаем шаблонный текст-заглушку
                    _replace_all_text(shape, " ")
                    cleared_count += 1

            # 2. Удаляем шаблонные изображения
            first_image_shape = None
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    if first_image_shape is None:
                        first_image_shape = shape
                    shapes_to_remove.append(shape)

            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
                deleted_images += 1

            # 3. Создаём текстовый плейсхолдер на месте первого изображения
            if image_placeholder and first_image_shape is not None:
                # Создаём текстбокс с теми же координатами
                txBox = slide.shapes.add_textbox(
                    first_image_shape.left, first_image_shape.top,
                    first_image_shape.width, first_image_shape.height
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"[{image_placeholder}]"
                run = p.runs[0]
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

        print(f"Заменено шейпов: {replacements_count}, очищено: {cleared_count}, удалено изображений: {deleted_images}")
        prs.save(str(output_path))
        print(f"Презентация сохранена: {output_path}")
        print(f"Слайдов: {len(prs.slides)}")


def _zip_directory(src_dir, output_path):
    """Запаковать директорию в ZIP (PPTX)."""
    with zipfile.ZipFile(str(output_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, dirs, files in os.walk(str(src_dir)):
            for fname in files:
                file_path = os.path.join(root, fname)
                arcname = os.path.relpath(file_path, str(src_dir))
                zf.write(file_path, arcname)



def _replace_all_text(shape, new_text):
    """Заменить весь текст, сохраняя формат первого run."""
    for para in shape.text_frame.paragraphs:
        if para.runs:
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""
            break
    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        p_elem = tf.paragraphs[-1]._p
        p_elem.getparent().remove(p_elem)


# === CLI commands ===

def list_slides(template_path=None):
    if template_path is None:
        template_path = TEMPLATE_PATH
    prs = Presentation(str(template_path))
    print(f"Шаблон: {template_path}")
    print(f"Всего слайдов: {len(prs.slides)}\n")
    for idx, slide in enumerate(prs.slides):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
        images = sum(1 for s in slide.shapes if s.shape_type == 13)
        print(f"Слайд {idx + 1}: {len(slide.shapes)} shapes, {images} images")
        for t in texts[:5]:
            print(f"  → {t[:100]}")
        print()


def _inspect_slide_data(slide, slide_num):
    """Вывести информацию о слайде (без загрузки файла)."""
    print(f"=== Слайд {slide_num} ===")
    print(f"Layout: {slide.slide_layout.name}")
    print(f"Shapes: {len(slide.shapes)}\n")
    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i}: type={shape.shape_type}, "
              f"pos=({shape.left}, {shape.top}), "
              f"size=({shape.width}x{shape.height})")
        if hasattr(shape, "text") and _normalize_text(shape.text):
            print(f"  Text: {_normalize_text(shape.text)[:200]}")
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font_info = []
                        if run.font.size:
                            font_info.append(f"size={run.font.size}")
                        if run.font.bold:
                            font_info.append("bold")
                        if run.font.color and run.font.color.rgb:
                            font_info.append(f"color={run.font.color.rgb}")
                        if font_info:
                            print(f"    Run: '{run.text[:80]}' [{', '.join(font_info)}]")
        if shape.shape_type == 13:
            print(f"  [IMAGE]")
        print()


def inspect_slides(slide_nums, template_path=None):
    """Осмотреть несколько слайдов за одну загрузку файла."""
    if template_path is None:
        template_path = TEMPLATE_PATH
    prs = Presentation(str(template_path))
    for slide_num in slide_nums:
        if slide_num < 1 or slide_num > len(prs.slides):
            print(f"Слайд {slide_num} не найден (всего {len(prs.slides)})\n")
            continue
        _inspect_slide_data(prs.slides[slide_num - 1], slide_num)


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    command = sys.argv[1]
    if command == "build":
        if len(sys.argv) < 4:
            print("Использование: pptx_builder.py build <config.json> <output.pptx>")
            sys.exit(1)
        with open(sys.argv[2]) as f:
            config = json.load(f)
        build_presentation(config, sys.argv[3])
    elif command == "list-slides":
        list_slides(sys.argv[2] if len(sys.argv) > 2 else None)
    elif command == "inspect":
        if len(sys.argv) < 3:
            print("Использование: pptx_builder.py inspect <slide_num> [slide_num2 ...]")
            sys.exit(1)
        inspect_slides([int(num) for num in sys.argv[2:]])
    else:
        print(f"Неизвестная команда: {command}")
        sys.exit(1)


if __name__ == "__main__":
    main()
